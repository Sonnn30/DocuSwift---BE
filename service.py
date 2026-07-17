import httpx
import openpyxl
from pptx import Presentation
import io
import pypdf
import docx
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings

# Proses indexing

# 1. extract text data dari file yang di upload
    # - karena jenis filenya berbeda-beda maka akan penyesuaian untuk tiap jenis file,
    #   seperti pada excel dan powerpoint
# 2. Split hasil data text yang sudah berhasil di extract menjadi chunk
# 3. Embed tiap chunk untuk menjadi vector data (jangan lupa update model untuk menyesuiakan dimensi model embedding)
# 4. Simpan tiap chunk yang sudah di embed ke dalam database

async def extract_file(file_url: str, file_type: str):
    async with httpx.AsyncClient() as client:
        response = await client.get(file_url)
 
    if file_type == "xlsx":
        wb = openpyxl.load_workbook(io.BytesIO(response.content))
        result = []
        for sheet_index, sheet_name in enumerate(wb.sheetnames, start=1):
            ws = wb[sheet_name]
            result.append({
                "sheet": sheet_index,
                "sheet_name": sheet_name,
                "data": [[cell.value for cell in row] for row in ws.iter_rows()]
            })
        return result
 
    if file_type == "pptx":
        prs = Presentation(io.BytesIO(response.content))
        result = []
        for slide_index, slide in enumerate(prs.slides, start=1):
            result.append({
                "slide": slide_index,
                "data": [shape.text for shape in slide.shapes if shape.has_text_frame]
            })
        return result
 
    if file_type == "pdf":
        pdf = pypdf.PdfReader(io.BytesIO(response.content))
        result = []
        for page_index, page in enumerate(pdf.pages, start=1):
            result.append({
                "page": page_index,
                "data": page.extract_text()
            })
        return result
 
    if file_type == "docx":
        doc = docx.Document(io.BytesIO(response.content))
        result = []
        page_index = 1
        current_page_text = ""
 
        for paragraph in doc.paragraphs:
            # cek apakah ada page break di paragraph ini
            if any(run.text == "" and "<w:lastRenderedPageBreak/>" in run._r.xml or "w:type" in run._r.xml and "page" in run._r.xml for run in paragraph.runs):
                result.append({"page": page_index, "data": current_page_text.strip()})
                page_index += 1
                current_page_text = ""
            else:
                current_page_text += paragraph.text + "\n"
 
        # simpan sisa teks halaman terakhir
        if current_page_text.strip():
            result.append({"page": page_index, "data": current_page_text.strip()})
 
        return result
 
    return response.text


def chunking(text: str):
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size = 400,
        chunk_overlap = 60
    )
    chunk = text_splitter.split_text(text)
    return chunk

embedding_model = HuggingFaceEmbeddings(
    model_name="paraphrase-multilingual-MiniLM-L12-v2"
)

def embedding(chunk: list):
    embeded = embedding_model.embed_documents(chunk)

    return embeded

def embedding_msg(msg: str):
    embeded = embedding_model.embed_query(msg)

    return embeded


def to_string(data):
    if isinstance(data, str):
        return data
    if isinstance(data, list):
        return " ".join([to_string(item) for item in data if item is not None])
    return str(data) if data is not None else ""
 

async def indexing(file_url: str, file_type: str):
    extracted = await extract_file(file_url, file_type)
 
    result = []
    for item in extracted:
        text = to_string(item["data"])
        chunks = chunking(text)
        embedded = embedding(chunks)
 
        for chunk, vector in zip(chunks, embedded):
            result.append({
                "metadata": {k: v for k, v in item.items() if k != "data"},
                "chunk": chunk,
                "vector": vector
            })
 
    return result